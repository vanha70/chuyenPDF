import streamlit as st
import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# 1. CẤU HÌNH & CSS (GIỮ NGUYÊN GIAO DIỆN ĐẸP)
# ==============================================================================
st.set_page_config(page_title="PDF to PowerPoint - Nguyễn Văn Hà", page_icon="⚡", layout="wide", initial_sidebar_state="collapsed")

st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&display=swap');
    .stApp { background-color: #020617; color: white; font-family: 'Inter', sans-serif; }
    header[data-testid="stHeader"] {display: none;}
    .header-container { display: flex; justify-content: space-between; align-items: center; padding: 10px 0px; border-bottom: 1px solid #1e293b; margin-bottom: 40px; }
    .logo-section { display: flex; align-items: center; gap: 15px; }
    .logo-icon { background: linear-gradient(135deg, #0ea5e9, #2563eb); color: white; width: 45px; height: 45px; border-radius: 10px; display: flex; align-items: center; justify-content: center; font-weight: bold; font-size: 24px; box-shadow: 0 0 15px rgba(14, 165, 233, 0.5); }
    .brand-name { font-size: 20px; font-weight: 800; letter-spacing: 1px; text-transform: uppercase; color: #ffffff; line-height: 1.2; }
    .brand-sub { font-size: 10px; color: #0ea5e9; font-weight: 600; letter-spacing: 1.5px; }
    .status-badge { background-color: rgba(34, 197, 94, 0.1); color: #22c55e; border: 1px solid #22c55e; padding: 5px 15px; border-radius: 20px; font-size: 11px; font-weight: bold; display: inline-flex; align-items: center; gap: 5px; }
    .dot { height: 8px; width: 8px; background-color: #22c55e; border-radius: 50%; display: inline-block; }
    .hero-title { text-align: center; font-size: 56px; font-weight: 900; margin-bottom: 10px; text-transform: uppercase; }
    .gradient-text { background: linear-gradient(to right, #fb923c, #fca5a5, #fff, #67e8f9); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }
    .hero-desc { text-align: center; color: #94a3b8; font-size: 18px; max-width: 700px; margin: 0 auto 60px auto; }
    div[data-testid="stFileUploader"] { border: 2px dashed #334155; border-radius: 15px; padding: 30px; background-color: rgba(30, 41, 59, 0.5); text-align: center; transition: all 0.3s ease; }
    div[data-testid="stFileUploader"]:hover { border-color: #f97316; background-color: rgba(249, 115, 22, 0.05); }
    div.stButton > button, div.stDownloadButton > button { width: 100%; background-color: #1e293b; color: #94a3b8; border: none; padding: 20px; font-size: 16px; font-weight: 800; border-radius: 12px; text-transform: uppercase; letter-spacing: 1px; transition: all 0.3s; height: 80px; }
    div.stButton > button:hover { background-color: #0ea5e9; color: white; box-shadow: 0 0 20px rgba(14, 165, 233, 0.4); }
    div.stDownloadButton > button { background-color: rgba(34, 197, 94, 0.2); color: #22c55e; border: 1px solid #22c55e; }
    div.stDownloadButton > button:hover { background-color: #22c55e; color: white; box-shadow: 0 0 20px rgba(34, 197, 94, 0.4); }
    .custom-card { background-color: #0f172a; border: 1px solid #1e293b; border-radius: 24px; padding: 40px; height: 100%; min-height: 350px; display: flex; flex-direction: column; justify-content: center; }
</style>
""", unsafe_allow_html=True)

# ==============================================================================
# 2. XỬ LÝ TEXT & CÔNG THỨC HÓA HỌC (CORE LOGIC)
# ==============================================================================

def clean_text(text):
    """Loại bỏ các tag và khoảng trắng thừa"""
    text = re.sub(r'\', '', text)
    return text.strip()

def format_chemical_text(paragraph, text, font_size=18, is_bold=False, color=None):
    """
    Hàm này cực kỳ quan trọng: Tự động phát hiện công thức hóa học 
    để in chỉ số trên (superscript) và chỉ số dưới (subscript).
    Ví dụ: H2SO4 -> 2, 4 xuống dưới; Cu2+ -> 2+ lên trên.
    """
    paragraph.clear() # Xóa text cũ nếu có
    p = paragraph
    
    # Regex để tách các phần: Chữ thường, số (sub), ion (super)
    # Logic đơn giản: Số đứng sau chữ cái -> Subscript. Dấu +, - đứng sau số -> Superscript
    
    # Tạm thời dùng logic tách từ đơn giản để xử lý Subscript cho số
    tokens = re.split(r'(\d+[+-]?|\s+)', text)
    
    for token in tokens:
        if not token: continue
        
        run = p.add_run()
        run.font.size = Pt(font_size)
        run.font.name = 'Arial'
        run.font.bold = is_bold
        if color:
            run.font.color.rgb = color
            
        # Kiểm tra nếu là công thức hóa học (Ví dụ: 2, 4 trong H2SO4)
        if re.match(r'^\d+$', token):
            # Nếu là số đứng riêng, thường là chỉ số dưới trong hóa học (trừ hệ số cân bằng)
            # Ở đây ta set mặc định là subscript cho đẹp với H2SO4
            run.text = token
            run.font.subscript = True
        elif re.match(r'^\d*[+-]$', token): # Ví dụ: 2+, +, -
            run.text = token
            run.font.superscript = True
        else:
            run.text = token

def parse_exam_content(full_content):
    """
    Hàm phân tích nội dung text thô thành danh sách câu hỏi có cấu trúc.
    Dựa trên cấu trúc file bạn cung cấp.
    """
    questions = []
    lines = full_content.split('\n')
    
    current_q = None
    state = "START" # START, QUESTION, OPTIONS
    
    # Regex phát hiện bắt đầu câu hỏi mới (Số đứng một mình hoặc dòng có số)
    # Trong file của bạn: 1 -> Dòng chỉ có số 1
    
    for line in lines:
        clean_line = clean_text(line)
        if not clean_line: continue
        
        # 1. Phát hiện số thứ tự câu hỏi (Vd: "1", "2", "28")
        if re.match(r'^\d+$', clean_line):
            if current_q: questions.append(current_q)
            current_q = {
                "id": clean_line,
                "content": "",
                "options": [],
                "type": "mcq" # Mặc định là trắc nghiệm
            }
            state = "QUESTION"
            continue
            
        # 2. Bỏ qua dòng chữ "CÂU HỎI" vô nghĩa
        if clean_line.upper() == "CÂU HỎI":
            continue
            
        # 3. Phát hiện đáp án A, B, C, D hoặc Đúng/Sai
        # Pattern: Bắt đầu bằng A, B, C, D nằm riêng hoặc "a.", "b."
        if re.match(r'^[A-D]$', clean_line) or re.match(r'^[a-d]\.', clean_line):
             if current_q:
                # Nếu gặp a. b. c. d. -> Chuyển sang dạng câu hỏi Đúng/Sai (Câu 19-22)
                if re.match(r'^[a-d]\.', clean_line):
                    current_q['type'] = "true_false"
                
                current_q['options'].append({"label": clean_line, "text": ""})
                state = "OPTIONS"
             continue

        # 4. Phát hiện dấu chọn đáp án đúng (✦) hoặc kết quả ĐÚNG/SAI/ĐÁP SỐ
        if "✦" in clean_line:
             if current_q and current_q['options']:
                 current_q['options'][-1]['is_correct'] = True
             continue
             
        if clean_line.startswith("ĐÚNG ✔") or clean_line.startswith("SAI ✘"):
             if current_q and current_q['options']:
                 current_q['options'][-1]['result'] = clean_line
             continue
        
        if clean_line.startswith("ĐÁP SỐ:"):
            if current_q:
                current_q['type'] = "short_ans"
                current_q['answer_text'] = clean_line
            continue

        # 5. Nạp nội dung
        if current_q:
            if state == "QUESTION":
                # Cộng dồn nội dung câu hỏi
                if "HỆ THỐNG GIÁO DỤC" not in clean_line: # Bỏ footer lẫn vào
                    current_q['content'] += clean_line + " "
            elif state == "OPTIONS":
                # Cộng dồn nội dung đáp án
                if current_q['options']:
                    if "HỆ THỐNG GIÁO DỤC" not in clean_line:
                        current_q['options'][-1]['text'] += clean_line + " "

    if current_q: questions.append(current_q)
    return questions

# ==============================================================================
# 3. TẠO SLIDE POWERPOINT (RENDER ENGINE)
# ==============================================================================

def create_pptx_file(questions):
    prs = Presentation()
    prs.slide_width = Inches(13.33) # Tỉ lệ 16:9
    prs.slide_height = Inches(7.5)

    # Màu sắc
    ORANGE = RGBColor(237, 125, 49)
    NAVY = RGBColor(0, 32, 96)
    GRAY = RGBColor(120, 120, 120)

    for q in questions:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank slide
        
        # --- 1. SỐ CÂU HỎI (Hộp cam bên trái) ---
        shape_num = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.8), Inches(0.8))
        shape_num.fill.solid()
        shape_num.fill.fore_color.rgb = ORANGE
        shape_num.line.color.rgb = ORANGE
        p_num = shape_num.text_frame.paragraphs[0]
        p_num.text = str(q['id'])
        p_num.font.size = Pt(36)
        p_num.font.bold = True
        p_num.alignment = PP_ALIGN.CENTER
        
        # --- 2. LABEL "CÂU HỎI" ---
        tb_lbl = slide.shapes.add_textbox(Inches(1.4), Inches(0.5), Inches(2), Inches(0.5))
        p_lbl = tb_lbl.text_frame.paragraphs[0]
        p_lbl.text = "CÂU HỎI"
        p_lbl.font.size = Pt(24)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = ORANGE

        # --- 3. NỘI DUNG CÂU HỎI ---
        tb_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1.5))
        tb_content.text_frame.word_wrap = True
        # Dùng hàm format để xử lý công thức hóa học trong câu hỏi
        format_chemical_text(tb_content.text_frame.paragraphs[0], q['content'], font_size=24, is_bold=True, color=NAVY)

        # --- 4. XỬ LÝ ĐÁP ÁN THEO LOẠI ---
        start_y = 3.5
        
        # TRƯỜNG HỢP A: TRẮC NGHIỆM (MCQ) - Chia 2 cột
        if q['type'] == 'mcq':
            col_coords = [(Inches(1.0), Inches(3.5)), (Inches(7.0), Inches(3.5)), 
                          (Inches(1.0), Inches(5.0)), (Inches(7.0), Inches(5.0))]
            
            for idx, opt in enumerate(q['options']):
                if idx >= 4: break
                left, top = col_coords[idx]
                
                # Chữ cái A, B, C, D
                labels = ["A", "B", "C", "D"]
                tb_opt_lbl = slide.shapes.add_textbox(left - Inches(0.5), top, Inches(0.5), Inches(0.5))
                p_opt_lbl = tb_opt_lbl.text_frame.paragraphs[0]
                p_opt_lbl.text = labels[idx]
                p_opt_lbl.font.size = Pt(24)
                p_opt_lbl.font.bold = True
                p_opt_lbl.font.color.rgb = ORANGE
                
                # Nội dung đáp án
                tb_opt_txt = slide.shapes.add_textbox(left, top, Inches(5.5), Inches(1.2))
                tb_opt_txt.text_frame.word_wrap = True
                format_chemical_text(tb_opt_txt.text_frame.paragraphs[0], opt['text'], font_size=20)
                
                # Nếu là đáp án đúng (có dấu ✦)
                if opt.get('is_correct'):
                    # Vẽ khung viền hoặc đổi màu để highlight
                    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - Inches(0.6), top - Inches(0.1), Inches(6), Inches(1.3))
                    rect.fill.background() # Trong suốt
                    rect.line.color.rgb = RGBColor(255, 0, 0)
                    rect.line.width = Pt(2)
                    # Đưa khung ra sau chữ
                    # (Python-pptx add theo thứ tự layer, nên add sau sẽ đè lên. 
                    # Ở đây ta chấp nhận vẽ đè lên text box một chút hoặc add rect trước text box)

        # TRƯỜNG HỢP B: ĐÚNG SAI (True/False) - List dọc
        elif q['type'] == 'true_false':
            for idx, opt in enumerate(q['options']):
                top = Inches(start_y + idx * 0.9)
                
                tb_row = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.8))
                p_row = tb_row.text_frame.paragraphs[0]
                
                full_text = f"{opt['label']} {opt['text']}"
                if opt.get('result'):
                    full_text += f"   [{opt['result']}]"
                
                format_chemical_text(p_row, full_text, font_size=20)

        # TRƯỜNG HỢP C: ĐIỀN ĐÁP ÁN (Short Ans)
        elif q['type'] == 'short_ans':
             tb_ans = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(10), Inches(1.0))
             p_ans = tb_ans.text_frame.paragraphs[0]
             p_ans.text = q.get('answer_text', '')
             p_ans.font.size = Pt(24)
             p_ans.font.bold = True
             p_ans.font.color.rgb = RGBColor(255, 0, 0)

        # --- 5. FOOTER ---
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.33), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(220, 220, 220)
        
        tb_footer = slide.shapes.add_textbox(Inches(0), Inches(7.1), Inches(13.33), Inches(0.4))
        p_footer = tb_footer.text_frame.paragraphs[0]
        p_footer.text = "HỆ THỐNG GIÁO DỤC HIỆN ĐẠI | BIÊN SOẠN: THẦY NGUYỄN VĂN HÀ"
        p_footer.font.size = Pt(12)
        p_footer.font.color.rgb = GRAY
        p_footer.alignment = PP_ALIGN.CENTER

    # Lưu file
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ==============================================================================
# 4. GIAO DIỆN STREAMLIT
# ==============================================================================

# HEADER
st.markdown("""
<div class="header-container">
    <div class="logo-section"><div class="logo-icon">H</div><div><div class="brand-name">NGUYỄN VĂN HÀ</div><div class="brand-sub">AI EDUCATION • DIGITAL TRANSFORMATION</div></div></div>
    <div class="contact-info"><div style="font-size: 10px; color: #64748b; margin-bottom: 2px;">HỖ TRỢ 24/7</div><div class="phone-number">0927.2222.05</div></div>
    <div class="status-badge"><span class="dot"></span> AI NODE ACTIVE</div>
</div>
""", unsafe_allow_html=True)

# HERO SECTION
st.markdown("""
<div style="margin-top: 50px;">
    <h1 class="hero-title"><span style="color: #f97316;">PDF</span> <span style="color: white;">TO</span> <span class="gradient-text">POWERPOINT</span> <span style="color: white;">SIÊU TỐC</span></h1>
    <p class="hero-desc">Hệ thống AI chuyên dụng giúp thầy cô chuyển đổi 100% học liệu sang PowerPoint tương tác chỉ với 1 cú nhấp chuột.</p>
</div>
""", unsafe_allow_html=True)

# MAIN UI
_, main_col, _ = st.columns([1, 8, 1])
with main_col:
    col1, col2 = st.columns(2, gap="large")

    # CỘT 1: NHẬP DỮ LIỆU
    with col1:
        st.markdown('<div class="custom-card"><div class="step-header" style="color:#f97316">BƯỚC 1: DỮ LIỆU</div>', unsafe_allow_html=True)
        # Vì bạn đã gửi nội dung text, ta sẽ để sẵn text đó vào đây để demo luôn
        # Hoặc cho phép upload file txt nếu muốn
        uploaded_file = st.file_uploader("Chọn file Text/Word đã convert", type=['txt', 'docx'])
        
        # Dữ liệu mẫu mặc định (Lấy từ file bạn gửi)
        default_content = """1
CÂU HỎI
Cấu trúc mạch vòng của carbohydrate nào sau đây không có nhóm -OH hemiacetal hoặc hemiketal?
A
Saccharose.
✦
B
Maltose.
C
Glucose.
D
Fructose.
2
CÂU HỎI
Carbohydrate nào sau đây kém tan trong nước lạnh nhưng tan được trong nước nóng tạo dung dịch keo, nhớt?
A
Cellulose.
B
Saccharose.
C
Tinh bột.
✦
D
Glucose.
10
CÂU HỎI
Khi pin Galvani Zn – Cu hoạt động thì
A
dòng electron chạy từ Cu sang Zn.
B
ở điện cực dương, cathode xảy ra quá trình oxi hóa Cu.
C
Zn đóng vai trò cực âm, Cu đóng vai trò cực dương.
✦
D
ở điện cực âm, anode xảy ra quá trình khử Zn.
23
CÂU HỎI
Hiện nay mạ điện được sử dụng rộng rãi trong thực tế... mạ Ag... I = 2A...
ĐÁP SỐ: 0,15
"""
        # Nếu chưa upload thì dùng text mẫu, nếu upload thì đọc file
        if uploaded_file:
            stringio = io.StringIO(uploaded_file.getvalue().decode("utf-8"))
            content_input = stringio.read()
            st.success(f"Đã đọc file: {uploaded_file.name}")
        else:
            st.info("Đang sử dụng dữ liệu mẫu (bạn có thể copy paste toàn bộ nội dung file vào đây)")
            content_input = st.text_area("Nội dung thô:", value=default_content, height=200)

        st.markdown('</div>', unsafe_allow_html=True)

    # CỘT 2: XỬ LÝ
    with col2:
        st.markdown('<div class="custom-card"><div class="step-header" style="color:#06b6d4">BƯỚC 2: XUẤT POWERPOINT</div>', unsafe_allow_html=True)
        
        if st.button("BẮT ĐẦU CHUYỂN ĐỔI"):
            import time
            with st.spinner("AI đang phân tích cú pháp & tạo công thức hóa học..."):
                time.sleep(1)
                
                # 1. Parse dữ liệu
                try:
                    questions_data = parse_exam_content(content_input)
                    st.toast(f"Đã tìm thấy {len(questions_data)} câu hỏi!")
                    
                    # 2. Tạo PPT
                    pptx_file = create_pptx_file(questions_data)
                    
                    # 3. Lưu Session
                    st.session_state.pptx_out = pptx_file
                    st.success("Xử lý thành công!")
                    
                except Exception as e:
                    st.error(f"Lỗi xử lý: {e}")

        if 'pptx_out' in st.session_state:
            st.download_button(
                label="📥 TẢI POWERPOINT HOÀN CHỈNH",
                data=st.session_state.pptx_out,
                file_name="Giao_An_Hoa_Hoc_AI.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        st.markdown('</div>', unsafe_allow_html=True)

st.markdown("<br><br>", unsafe_allow_html=True)
